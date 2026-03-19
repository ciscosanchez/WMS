from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DECK_TITLE = "Armstrong WMS"
DECK_SUBTITLE = "Cleaned presentation of the current document set"

SLIDES = [
    {
        "type": "title",
        "title": DECK_TITLE,
        "subtitle": DECK_SUBTITLE,
    },
    {
        "title": "Architecture",
        "subtitle": "System Overview",
        "sections": [
            ("Application stack", [
                "Next.js App Router, Prisma ORM, NextAuth v5, and MinIO/S3 storage.",
                "Middleware layer handles tenant resolution, authentication, and RBAC.",
                "PostgreSQL 16 stores shared platform data and tenant-specific WMS data.",
            ]),
            ("Data isolation", [
                "Schema-per-tenant model: public schema plus tenant-specific schemas such as tenant_acme.",
                "Chosen over shared tables for stronger isolation without database-per-tenant overhead.",
            ]),
        ],
        "source": "docs/architecture.md",
    },
    {
        "title": "Architecture",
        "subtitle": "Operational Design",
        "sections": [
            ("Dual-mode operations", [
                "Shared core: inventory, products, warehouse, clients, audit log, and sequences.",
                "Freight/3PL mode: inbound ASNs, BOL processing, container receiving, discrepancies, inspections.",
                "Fulfillment mode: orders, sales channels, pick tasks, packing, outbound shipments, carrier rate shopping.",
            ]),
            ("Core patterns", [
                "Immutable inventory transaction ledger records receive, putaway, move, pick, and adjust events.",
                "Request lifecycle flows through middleware, auth, tenant resolution, business logic, audit, and cache revalidation.",
            ]),
        ],
        "source": "docs/architecture.md",
    },
    {
        "title": "Ecosystem",
        "subtitle": "System Inventory",
        "sections": [
            ("Existing systems", [
                "NetSuite ERP for financials, billing, AR/AP, customer records, and contracts.",
                "DispatchPro TMS for load management, carrier dispatch, and tracking.",
            ]),
            ("Building / planned systems", [
                "Armstrong WMS for receiving, inventory, fulfillment, and warehouse operations.",
                "Operator mobile app, client portal, marketplace connectors, carrier integrations, and document intelligence.",
            ]),
            ("Identified gaps", [
                "Sales CRM, yard management, and billing engine.",
            ]),
        ],
        "source": "docs/ecosystem.md",
    },
    {
        "title": "Ecosystem",
        "subtitle": "Application Structure And Integration Flows",
        "sections": [
            ("Application structure", [
                "Single codebase with route groups for tenant WMS, operator workflows, client portal, and platform admin.",
                "Operator app and client portal are part of the same application, not separate codebases.",
            ]),
            ("Critical integration flows", [
                "WMS billable events to NetSuite.",
                "WMS ship requests and tracking with DispatchPro.",
                "Orders, inventory, and tracking with marketplaces and EDI.",
                "Label requests and tracking with carrier APIs.",
                "Real-time inventory, order, shipment, and billing visibility through the client portal.",
            ]),
        ],
        "source": "docs/ecosystem.md",
    },
    {
        "title": "Identification To Cash",
        "subtitle": "Phase Structure",
        "sections": [
            ("Phase 1: Identify & sell", [
                "Lead capture, qualification, proposal, contract, and onboarding.",
                "CRM owns lead-to-opportunity flow; NetSuite stores customer and rate-card terms; WMS sets up client and warehouse context.",
            ]),
            ("Phase 2: Receive & store", [
                "DispatchPro manages carrier scheduling and tracking.",
                "WMS creates ASN and receiving context; operator workflow executes unload, receive, and putaway.",
            ]),
            ("Phase 3: Store & manage", [
                "Storage billing, cycle counts, replenishment, alerts, and client inventory visibility.",
            ]),
        ],
        "source": "docs/flow-identification-to-cash.md",
    },
    {
        "title": "Identification To Cash",
        "subtitle": "Fulfill, Bill, And Collect",
        "sections": [
            ("Phase 4: Fulfill & ship", [
                "Orders arrive from client portal, marketplaces, EDI 940, API, or manual entry.",
                "WMS handles allocation, batching, pick tasks, packing, shipping decisions, and tracking.",
            ]),
            ("Phase 5: Bill & collect", [
                "Warehouse events become billable events that move into NetSuite for invoicing and collection.",
                "Examples across the document include receiving, storage, pick/pack, value-add services, and shipping charges.",
            ]),
        ],
        "source": "docs/flow-identification-to-cash.md",
    },
    {
        "title": "Receiving Flow",
        "subtitle": "Status Workflow",
        "sections": [
            ("Statuses", [
                "Draft -> Expected -> Arrived -> Receiving -> Inspection -> Completed.",
                "Any status can transition to Cancelled.",
            ]),
            ("Purpose", [
                "ASN-based freight receiving workflow for expected inbound shipments.",
                "Registers expected freight before arrival and finalizes inventory after warehouse processing.",
            ]),
        ],
        "source": "docs/flow-receiving.md",
    },
    {
        "title": "Receiving Flow",
        "subtitle": "Step-By-Step Workflow",
        "sections": [
            ("Setup", [
                "Create ASN draft with client, carrier, BOL, tracking, PO number, and expected date.",
                "Add line items with product, expected quantity, UOM, and lot number when applicable.",
                "Mark shipment expected and notify the warehouse.",
            ]),
            ("Execution", [
                "Mark arrival at the dock.",
                "Receive line items with actual quantity, condition, bin, lot/serial, and notes.",
                "Log discrepancies for shortage, overage, or damage.",
                "Run optional inspection checklists.",
                "Complete receiving and create inventory records plus receive transactions.",
            ]),
        ],
        "source": "docs/flow-receiving.md",
    },
    {
        "title": "Fulfillment Flow",
        "subtitle": "Order Status Workflow",
        "sections": [
            ("Statuses", [
                "Pending -> Awaiting Fulfillment -> Allocated -> Picking -> Picked -> Packing -> Packed -> Shipped -> Delivered.",
                "Exception statuses include Cancelled, On Hold, and Backordered.",
            ]),
            ("Purpose", [
                "Outbound workflow for orders from Shopify, Amazon, Walmart, manual entry, and API.",
                "Orders are ingested, allocated, picked, packed, and shipped with tracking.",
            ]),
        ],
        "source": "docs/flow-fulfillment.md",
    },
    {
        "title": "Fulfillment Flow",
        "subtitle": "Step-By-Step Workflow",
        "sections": [
            ("Order to pick", [
                "Ingest order with client, ship-to, line items, priority, and ship-by date.",
                "Allocate inventory and create allocation transactions.",
                "Generate pick tasks using single-order, batch, wave, or zone methods.",
            ]),
            ("Pick to ship", [
                "Pick execution confirms bin, product, quantity, and short-pick handling.",
                "Packing records package choice, dimensions, weight, and value-added services.",
                "Shipping rate-shops carriers, generates labels, stores tracking, and confirms delivery.",
            ]),
        ],
        "source": "docs/flow-fulfillment.md",
    },
    {
        "title": "Document Intelligence",
        "subtitle": "Scan -> Extract -> Verify -> Post",
        "sections": [
            ("Problem addressed", [
                "Manual keying of BOLs, packing lists, invoices, PODs, purchase orders, customs documents, and tally sheets.",
                "Current pain points: slow processing, error risk, labor cost, and dock bottlenecks.",
            ]),
            ("Proposed workflow", [
                "Capture from phone camera, scanner, email attachment, upload, EDI inbound, or PDF.",
                "AI/OCR extracts structured fields with confidence scores.",
                "Human review confirms low-confidence fields before posting into WMS workflows.",
            ]),
        ],
        "source": "docs/document-intelligence.md",
    },
    {
        "title": "Competitive Analysis",
        "subtitle": "Alternatives Reviewed",
        "sections": [
            ("Compared platforms", [
                "Logiwa IO",
                "Oracle NetSuite WMS",
                "Magaya",
                "3PL Warehouse Manager (Extensiv)",
            ]),
            ("Comparison dimensions", [
                "Primary market, domain focus, multi-tenancy, architecture, receiving, inventory, integrations, pricing, and billing.",
            ]),
            ("Document takeaway", [
                "Armstrong is positioned as combining freight-oriented receiving depth with fulfillment breadth.",
            ]),
        ],
        "source": "docs/competitive-analysis.md",
    },
    {
        "title": "Deployment",
        "subtitle": "Local And AWS Production",
        "sections": [
            ("Local development", [
                "Docker Compose for PostgreSQL and MinIO.",
                "Next.js dev server runs on the host and connects to local containers.",
            ]),
            ("AWS production", [
                "App Runner or ECS Fargate for the application.",
                "RDS PostgreSQL 16 for database, S3 for file storage, and Route 53 plus CloudFront for DNS and SSL.",
            ]),
            ("Environment model", [
                "No code changes required between local and AWS; switch is driven by environment variables.",
            ]),
        ],
        "source": "docs/deployment.md",
    },
    {
        "title": "Source Documents",
        "subtitle": "Included In This Presentation",
        "sections": [
            ("Documents", [
                "docs/architecture.md",
                "docs/ecosystem.md",
                "docs/flow-identification-to-cash.md",
                "docs/flow-receiving.md",
                "docs/flow-fulfillment.md",
                "docs/document-intelligence.md",
                "docs/competitive-analysis.md",
                "docs/deployment.md",
            ]),
        ],
        "source": "",
    },
]


BG = RGBColor(249, 247, 242)
TEXT = RGBColor(34, 43, 52)
MUTED = RGBColor(100, 107, 114)
NAVY = RGBColor(27, 51, 70)
ACCENT = RGBColor(150, 94, 61)
BORDER = RGBColor(218, 221, 224)
WHITE = RGBColor(255, 255, 255)


def style_run(run, size, color, bold=False):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()


def add_page(slide, n):
    box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.4), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    style_run(r, 9, MUTED)


def add_title(prs, n, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tag = slide.shapes.add_textbox(Inches(0.9), Inches(0.9), Inches(2.0), Inches(0.3))
    p = tag.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "DOCUMENT PRESENTATION"
    style_run(r, 11, ACCENT, True)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8.0), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    style_run(r, 28, TEXT, True)

    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(2.35), Inches(9.5), Inches(0.7))
    subtitle.text_frame.word_wrap = True
    p = subtitle.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data["subtitle"]
    style_run(r, 16, MUTED)

    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.4), Inches(6.6), Inches(1.2)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = WHITE
    note.line.color.rgb = BORDER
    tf = note.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "This deck keeps the source documents intact in presentation format rather than reframing them as a sales pitch."
    style_run(r, 13, TEXT)

    add_page(slide, n)


def add_content(prs, n, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(0.72), Inches(10.0), Inches(0.45))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    style_run(r, 23, TEXT, True)

    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(1.22), Inches(10.0), Inches(0.35))
    p = subtitle.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data["subtitle"]
    style_run(r, 13, MUTED, True)

    left = 0.9
    top = 1.8
    width = 5.45
    height = 2.2
    gap = 0.35

    sections = data["sections"]
    for idx, (heading, bullets) in enumerate(sections[:4]):
        col = idx % 2
        row = idx // 2
        x = Inches(left + col * (width + gap))
        y = Inches(top + row * (height + 0.35))
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.12)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        style_run(r, 12, ACCENT, True)

        for bullet in bullets:
            p = tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = bullet
            style_run(r, 12.5, TEXT)

    if data.get("source"):
        foot = slide.shapes.add_textbox(Inches(0.95), Inches(6.95), Inches(4.2), Inches(0.2))
        p = foot.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"Source: {data['source']}"
        style_run(r, 9, MUTED)

    add_page(slide, n)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Armstrong WMS Cleaned Document Presentation"

    for i, slide_data in enumerate(SLIDES, start=1):
        if slide_data.get("type") == "title":
            add_title(prs, i, slide_data)
        else:
            add_content(prs, i, slide_data)

    out = "docs/Armstrong-WMS-Cleaned-Docs.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
